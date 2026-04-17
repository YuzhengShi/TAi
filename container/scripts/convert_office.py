#!/usr/bin/env python3
"""Convert Office files (.docx, .xlsx, .pptx) to readable plain text.

Usage: python3 /opt/scripts/convert_office.py <input_file>

Outputs extracted text to stdout. Exits non-zero on failure.
"""
import sys
import os


def convert_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(' | '.join(cells))
        if rows:
            parts.append('\n'.join(rows))
    return '\n'.join(parts)


def convert_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"=== Sheet: {sheet} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            if any(cells):
                parts.append(' | '.join(cells))
    wb.close()
    return '\n'.join(parts)


def convert_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(' | '.join(cells))
    return '\n'.join(parts)


def main():
    if len(sys.argv) != 2:
        print(f"Usage: {sys.argv[0]} <input_file>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.isfile(path):
        print(f"File not found: {path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    converters = {
        '.docx': convert_docx,
        '.xlsx': convert_xlsx,
        '.pptx': convert_pptx,
        '.doc': convert_docx,
        '.xls': convert_xlsx,
        '.ppt': convert_pptx,
    }

    converter = converters.get(ext)
    if not converter:
        print(f"Unsupported format: {ext}", file=sys.stderr)
        sys.exit(1)

    try:
        text = converter(path)
        if text:
            print(text)
        else:
            print(f"[No text content extracted from {os.path.basename(path)}]", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(f"Conversion failed: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
